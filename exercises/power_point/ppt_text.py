# ppt_text.py   29-Jul-2021  crs from from doc
"""
# https://python-pptx.readthedocs.io/en/latest
/user/quickstart.html
"""
import sys

from pptx import Presentation

args = sys.argv[1:]
if len(args) > 0:
    path_to_presentation = args[0]
else:
    path_to_presentation = input("ppt path:")


prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

slide_no = 0
for slide in prs.slides:
    slide_no += 1
    for shape in slide.shapes:
        print("Slide:", slide_no)
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print("    " + run.text)
