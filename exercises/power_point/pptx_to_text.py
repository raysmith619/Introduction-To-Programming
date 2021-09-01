from pptx import Presentation
import glob

for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    slide_no = 0
    for slide in prs.slides:
        slide_no += 1
        print(f"\nSide {slide_no}", end=" ")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
